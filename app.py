from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pytube import YouTube
import requests
from bs4 import BeautifulSoup
import os
from PySimpleGUI import PySimpleGUI as sg
from pathlib import Path


def formatar_inputs(entrada):
    entrada.replace(' ', '-')
    entrada.lower()
    entrada.strip()
    return entrada


def encontrar_letra(artista, nome_louvor):
    link = f"https://www.letras.com/{formatar_inputs(artista)}/{formatar_inputs(nome_louvor)}/"

    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/118.0.0.0 Safari/537.36"}
    requisicao = requests.get(link, headers=headers)

    site = BeautifulSoup(requisicao.text, "html.parser")

    letra_musica = str(site.find("div", class_="lyric-original"))
    letra_musica = letra_musica.replace('<br/>', '\n')
    letra_musica = letra_musica.replace('<p>', '\n')
    letra_musica = letra_musica.replace('</p>', '\n')
    letra_musica = letra_musica.replace('<div class="lyric-original">', '')
    letra_musica = letra_musica.replace('</div>', '')

    return letra_musica


def novo_slide(prss):
    slide_layout = prss.slide_layouts[6]  # 5 é o índice para o layout de título em preto
    slide = prss.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Cor preta
    return slide


def criar_capa(nome_louvor, nome_cantor, prss):
    # Adicione a capa usando a função criarCapa
    slide = novo_slide(prs)

    # Configure o fundo preto
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Cor preta

    # Adicione o nome da música em maiúsculas com fonte Arial e tamanho 48
    musica = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
    text_frame = musica.text_frame
    prgh = text_frame.add_paragraph()
    prgh.text = nome_louvor.strip()
    prgh.font.size = Pt(48)
    prgh.font.bold = True
    prgh.font.name = 'Arial'
    prgh.alignment = PP_ALIGN.CENTER
    prgh.font.color.rgb = RGBColor(255, 255, 255)  # Cor branca

    # Adicione o nome do compositor com a mesma fonte, mas tamanho 24
    compositor = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
    text_frame = compositor.text_frame
    prgh = text_frame.add_paragraph()
    prgh.text = nome_cantor.strip()
    prgh.font.size = Pt(24)
    prgh.font.bold = True
    prgh.font.name = 'Arial'
    prgh.alignment = PP_ALIGN.CENTER
    prgh.font.color.rgb = RGBColor(255, 255, 255)  # Cor branca

    # Centraliza verticalmente
    altura_slide = prs.slide_height
    altura_total = sum([paragrafo.space_after for paragrafo in text_frame.paragraphs if paragrafo.space_after])
    posicao_vertical = (altura_slide - altura_total) / 2 if altura_total else 0
    for paragrafo in text_frame.paragraphs:
        if paragrafo.space_after:
            paragrafo.space_after = Pt(0)  # Remova o espaço após os parágrafos
        if paragrafo.space_before:
            paragrafo.space_before = Pt(0)  # Remova o espaço antes dos parágrafos
        paragrafo.top = Inches(posicao_vertical)

    return prss


def baixar_louvor(url_louvor, titulo):
    # Baixar louvor
    try:
        yt = YouTube(url_louvor)
        yt.title = titulo
        audio_stream = yt.streams.filter(only_audio=True, file_extension='mp4').first()
        audio_stream.download(output_path='louvores')
        print('Download da música concluída!')

    except:
        print('Erro ao baixar o louvor, URL inválida!')


# Layout
sg.theme("Reddit")
layout = [
    [sg.Text("Link do Youtube:"), sg.Input(key="download_musica")],
    [sg.Text("Número da música:"), sg.Input(key="numero")],
    [sg.Text("Nome da música:"), sg.Input(key="louvor")],
    [sg.Text("Nome do cantor:"), sg.Input(key="cantor")],
    [sg.Button("Iniciar")]
]
# Janela
janela = sg.Window("Easy Lyrics", layout)
# Ler os eventos
while True:
    eventos, valores = janela.read()
    if eventos == sg.WINDOW_CLOSED:
        break

    if eventos == "Iniciar":

        # nome_arquivos = input("Digite o nome para os arquivos: (Ex: 1 nome da música - nome do cantor): ")
        nome_arquivo = f'{valores["numero"]} {valores["louvor"]} - {valores["cantor"]}'
        nome_musica = valores["louvor"]
        nome_compositor = valores["cantor"]

        print()

        with open('letra_musica.txt', 'w', encoding='utf-8') as file:
            file.write(encontrar_letra(nome_compositor, nome_musica))
            print('Escrevendo a letra da música no arquivo...')

        # Ler a letra da música de um arquivo
        with open('letra_musica.txt', 'r', encoding='utf-8') as file:
            estrofes = file.read().split('\n\n')  # Assumindo que as estrofes estão separadas por duas quebras de linha
            file.close()

        # Crie uma apresentação do PowerPoint em resolução Full HD (1920x1080)
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # 1920 pixels / 144 pixels por polegada
        prs.slide_height = Inches(7.5)  # 1080 pixels / 144 pixels por polegada

        # Configuração da fonte
        font_name = 'Arial'
        font_size = Pt(48)
        font_color = RGBColor(255, 255, 255)  # Branco

        # Chamando a função criarCapa no início do código
        criar_capa(nome_musica, nome_compositor, prs)
        baixar_louvor(valores["download_musica"], nome_arquivo)

        for estrofe in estrofes:
            try:
                if estrofe == "None":
                    print(f'A letra da música não foi encontrada no banco de dados. Faça o arquivo "{nome_arquivo}" à mão')
            except:
                os.remove(f'{nome_arquivo}.pptx')
                os.remove(f'{nome_arquivo}.mp4')

            slide = novo_slide(prs)

            # Adicione o texto da estrofe ao slide
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12.5)  # Largura do slide - margens
            height = Inches(4)  # Altura do slide - margens
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame

            p = tf.add_paragraph()
            p.text = estrofe
            p.font.name = font_name
            p.font.size = font_size
            p.font.bold = True  # Negrito
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(24)  # Espaço após o parágrafo em pontos
            p.font.color.rgb = font_color

        # Salve a apresentação em resolução
        output_path = f'louvores/{nome_arquivo}.pptx'
        prs.save(output_path)

        os.remove('letra_musica.txt')

        print("FIM!")

